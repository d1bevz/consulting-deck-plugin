"""Deck Assembly — Combine slide PNGs into PDF or PPTX."""

from pathlib import Path

from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from pptx import Presentation
from pptx.util import Inches


SLIDE_WIDTH = 13.333  # inches (16:9 at 96dpi standard)
SLIDE_HEIGHT = 7.5


def assemble_pdf(image_paths: list[str], output_path: str) -> str:
    """Combine PNG images into a single PDF, one image per page."""
    if not image_paths:
        raise ValueError("No images provided for assembly")

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    page_width = SLIDE_WIDTH * inch
    page_height = SLIDE_HEIGHT * inch

    c = canvas.Canvas(str(out), pagesize=(page_width, page_height))

    for img_path in image_paths:
        c.drawImage(img_path, 0, 0, width=page_width, height=page_height,
                     preserveAspectRatio=True, anchor="c")
        c.showPage()

    c.save()
    return str(out)


def assemble_pptx(image_paths: list[str], output_path: str) -> str:
    """Combine PNG images into a PPTX, one image per slide."""
    if not image_paths:
        raise ValueError("No images provided for assembly")

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_WIDTH)
    prs.slide_height = Inches(SLIDE_HEIGHT)

    blank_layout = prs.slide_layouts[6]  # Blank layout

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path, Inches(0), Inches(0),
            width=Inches(SLIDE_WIDTH), height=Inches(SLIDE_HEIGHT),
        )

    prs.save(str(out))
    return str(out)


if __name__ == "__main__":
    import sys
    import glob

    args = sys.argv[1:]
    input_dir = args[0] if args else "output"
    fmt = args[1] if len(args) > 1 else "pdf"

    images = sorted(glob.glob(f"{input_dir}/slide_*.png"))
    if not images:
        images = sorted(glob.glob(f"{input_dir}/demo_*.png"))

    if not images:
        print(f"No slide images found in {input_dir}/")
        sys.exit(1)

    output = f"output/deck.{fmt}"
    if fmt == "pptx":
        assemble_pptx(images, output)
    else:
        assemble_pdf(images, output)

    print(f"Assembled {len(images)} slides into {output}")
